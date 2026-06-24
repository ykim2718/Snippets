"""
y, Python import image to pptx.py, 2017.5.11

"""


def import_image_stream_to_pptx():

    import numpy as np
    import matplotlib.pyplot as plt
    from pptx import Presentation
    from io import BytesIO
    from pptx.util import Cm

    prs = Presentation()
    title_slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(title_slide_layout)

    x = np.arange(1, 101)
    y = 10 + 2 * x + np.random.normal(0, 60, 100)
    plt.plot(x, y, 'o')

    image_stream = BytesIO()
    plt.savefig(image_stream, format='png')
    image_stream.read()
    image_stream.seek(0)

    # pic = slide.shapes.add_picture(image_file=image_stream, left=Cm(1), top=Cm(1))
    image_stream.close()
    prs.save('sample image 1.pptx')


def import_image_file_to_pptx():

    import numpy as np
    import matplotlib.pyplot as plt
    from pptx import Presentation
    from io import StringIO, BytesIO
    from pptx.util import Cm

    prs = Presentation()
    title_slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(title_slide_layout)

    x = np.arange(1, 101)
    y = 10 + 2 * x + np.random.normal(0, 60, 100)
    plt.plot(x, y, 'o')

    image_file = 'sample image 2.png'
    plt.savefig(image_file)

    pic = slide.shapes.add_picture(image_file=image_file, left=Cm(1), top=Cm(1))
    prs.save('sample image 2.pptx')


if __name__ == '__main__':

    import_image_stream_to_pptx()
    import_image_file_to_pptx()


